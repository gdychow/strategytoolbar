"""Stateless PPTX -> per-slide thumbnails/extracts conversion sidecar.

No database access, no persistent volume mount, and no published port
outside the Docker Compose network (see ../docker-compose.yml's `render`
service) -- this container has nothing to protect and nothing to exfiltrate
even if compromised. Every request works entirely inside a fresh temp
directory that's deleted before the response is returned; nothing survives
past one request.

POST /convert with a multipart 'file' field (a .pptx/.potx) returns a zip:
  manifest.json  -- {"slides": [{index, title, thumbnail, insertMode,
                     reconstructSpec, pptx}, ...]}
  slide-N.png    -- whitespace-trimmed (with a small margin re-added, see
                     THUMBNAIL_MARGIN_PX) PNG thumbnail for slide N (1-based)
  slide-N.pptx   -- present only for insertMode "file": that slide alone,
                     as its own single-slide presentation

Pipeline verified directly against a real multi-shape .pptx before this was
written: LibreOffice (--convert-to pdf) -> pdftoppm (per-slide PNG) ->
ImageMagick (-trim +repage, then a small white border added back) ->
python-pptx (title hint, single-slide
extraction, and reconstruct/file classification -- the same technique and
the same classify_shape_tree/extract_reconstruct_spec logic
scripts/slice-catalog-source.py already uses for the offline bulk-seed
pipeline, ported here so admin bulk uploads get the same one-click
"reconstruct" fidelity for simple shapes instead of always falling back to
the heavier temp-slide/copy-paste "file" flow).
"""

import colorsys
import io
import json
import os
import re
import shutil
import subprocess
import tempfile
import zipfile

from lxml import etree
from flask import Flask, jsonify, request, send_file
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE, MSO_LINE_DASH_STYLE, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, MSO_UNDERLINE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 100 * 1024 * 1024  # real templates carry embedded fonts/images

SUBPROCESS_TIMEOUT_SECONDS = 120
THUMBNAIL_DPI = 150
# `-trim` crops to the exact non-background bounding box with zero
# breathing room, which made every thumbnail look uncomfortably tight
# against its own content -- this re-adds a small white margin on all
# sides afterward so there's a visible sense of the graphic sitting on a
# background rather than filling the frame edge-to-edge.
THUMBNAIL_MARGIN_PX = 12
THINK_CELL_MARKER = "think-cell"


class ConversionError(Exception):
    def __init__(self, message, status=400):
        super().__init__(message)
        self.message = message
        self.status = status


@app.errorhandler(ConversionError)
def handle_conversion_error(err):
    return jsonify({"error": err.message}), err.status


@app.errorhandler(413)
def handle_too_large(_err):
    return jsonify({"error": "File too large."}), 413


@app.get("/health")
def health():
    return jsonify({"ok": True})


def run(cmd):
    try:
        result = subprocess.run(cmd, capture_output=True, timeout=SUBPROCESS_TIMEOUT_SECONDS)
    except subprocess.TimeoutExpired:
        raise ConversionError(f"{cmd[0]} timed out.", status=422)
    if result.returncode != 0:
        stderr = result.stderr.decode("utf-8", "replace")[:500]
        raise ConversionError(f"{cmd[0]} failed: {stderr}", status=422)
    return result


def extract_title_hint(slide):
    """Prefers a real Title/Center-Title placeholder; falls back to the
    first non-empty text run on the slide. Same preference order already
    established in scripts/slice-catalog-source.py."""
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type in (
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
        ):
            text = shape.text_frame.text.strip() if shape.has_text_frame else ""
            if text:
                return text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text[:80]
    return None


def slice_single_slide(source_path, slide_index, out_path):
    """Reloads the source presentation fresh and deletes every slide except
    slide_index -- python-pptx has no clone operation, so mutating one
    shared Presentation object across iterations would progressively delete
    slides out from under later iterations."""
    prs = Presentation(source_path)
    slide_id_list = prs.slides._sldIdLst
    slide_ids = list(slide_id_list)
    for i, slide_id in enumerate(slide_ids):
        if i != slide_index:
            slide_id_list.remove(slide_id)
    prs.save(out_path)


# ---------------------------------------------------------------------------
# reconstruct-vs-file classification, ported directly from
# scripts/slice-catalog-source.py -- that script's own module docstring and
# per-function comments carry the full "why" (which OOXML shape features
# have no PowerPoint JS API equivalent, why an unmapped preset falls back to
# 'file' instead of guessing, etc.); this is a straight port, not a
# reinterpretation, so consult that file for the reasoning if it's not
# repeated verbatim here.
# ---------------------------------------------------------------------------


def is_think_cell_placeholder(shape) -> bool:
    return THINK_CELL_MARKER in shape._element.xml


def is_unused_layout_placeholder(shape) -> bool:
    """A layout placeholder (Title, Subtitle, Content, etc.) the admin
    never actually clicked into still shows up in slide.shapes as a real
    shape -- confirmed directly: a bare 'Title Slide'-layout slide with
    only one real shape added still has 3 shapes total (2 empty
    placeholders + the real one). An untouched placeholder carries no
    explicit <a:prstGeom> at all (its geometry is purely inherited from
    the layout, invisibly), so classify_shape_tree would otherwise
    misclassify the whole slide as 'file' purely because of an invisible,
    content-free placeholder sitting on it -- this is exactly why
    Slide.exportAsBase64()'s "capture the whole current slide" (used by
    the task pane's "Add to Library"/Edit flows) could classify
    differently than the Bulk Upload pipeline's pre-isolated single-shape
    source files for what is visually the identical piece of content.
    Filtered out the same way is_think_cell_placeholder already filters
    that other kind of invisible-but-present artifact. A placeholder the
    admin DID type real text into is a genuine content shape and is not
    filtered -- only a placeholder with no text at all is treated as
    unused."""
    if not shape.is_placeholder:
        return False
    return not (shape.has_text_frame and shape.text_frame.text.strip())


def emu_to_pt(value):
    return round(Emu(value).pt, 2) if value is not None else None


def has_cust_geom(element) -> bool:
    return element.find(".//" + qn("a:custGeom")) is not None


def has_picture(element) -> bool:
    return element.tag == qn("p:pic") or element.find(".//" + qn("p:pic")) is not None


def has_other_graphic_frame(element) -> bool:
    return element.tag == qn("p:graphicFrame") or element.find(".//" + qn("p:graphicFrame")) is not None


def has_custom_bullet_char(element) -> bool:
    return element.find(".//" + qn("a:buChar")) is not None


# Maps an OOXML <a:prstGeom prst="..."> value to the exact string
# PowerPoint.GeometricShapeType's TypeScript enum uses. NOT a mechanical
# capitalization -- a deliberately conservative, hand-checked subset against
# the real enum list (see slice-catalog-source.py). A prst not in this table
# routes to 'file' mode rather than risk addGeometricShape silently
# accepting a wrong-but-validly-spelled value.
PRST_TO_GEOMETRIC_SHAPE_TYPE = {
    "rect": "Rectangle",
    "ellipse": "Ellipse",
    "roundRect": "RoundRectangle",
    "triangle": "Triangle",
    "rtTriangle": "RightTriangle",
    "diamond": "Diamond",
    "parallelogram": "Parallelogram",
    "trapezoid": "Trapezoid",
    "pentagon": "Pentagon",
    "hexagon": "Hexagon",
    "heptagon": "Heptagon",
    "octagon": "Octagon",
    "decagon": "Decagon",
    "dodecagon": "Dodecagon",
    "star4": "Star4",
    "star5": "Star5",
    "star6": "Star6",
    "star7": "Star7",
    "star8": "Star8",
    "star10": "Star10",
    "star12": "Star12",
    "star16": "Star16",
    "star24": "Star24",
    "star32": "Star32",
    "chevron": "Chevron",
    "donut": "Donut",
    "heart": "Heart",
    "sun": "Sun",
    "moon": "Moon",
    "cube": "Cube",
    "can": "Can",
    "cloud": "Cloud",
    "wave": "Wave",
    "plus": "Plus",
    "arc": "Arc",
    "chord": "Chord",
    "bevel": "Bevel",
    "frame": "Frame",
    "corner": "Corner",
    "pie": "Pie",
    "blockArc": "BlockArc",
    "smileyFace": "SmileyFace",
    "lightningBolt": "LightningBolt",
    "plaque": "Plaque",
    "teardrop": "Teardrop",
    "homePlate": "HomePlate",
    "rightArrow": "RightArrow",
    "leftArrow": "LeftArrow",
    "upArrow": "UpArrow",
    "downArrow": "DownArrow",
    "leftRightArrow": "LeftRightArrow",
    "upDownArrow": "UpDownArrow",
    "leftBracket": "LeftBracket",
    "rightBracket": "RightBracket",
    "leftBrace": "LeftBrace",
    "rightBrace": "RightBrace",
    "wedgeRectCallout": "WedgeRectCallout",
    # Deliberately NOT mapped: prst="line" has no GeometricShapeType
    # equivalent -- a straight line needs PowerPoint.Slide.addLine, a
    # different API the reconstruct path doesn't call.
}


# a:pPr algn -> PowerPoint.ParagraphHorizontalAlignment's exact string
# values (confirmed against @types/office-js: Left/Center/Right/Justify/
# JustifyLow/Distributed/ThaiDistributed). PP_ALIGN.to_xml() round-trips a
# python-pptx alignment member straight back to its raw OOXML value -- same
# technique already used for MSO_THEME_COLOR below -- so this table only
# maps that raw value to Office.js's naming, it doesn't re-derive it.
ALGN_TO_OFFICEJS_ALIGNMENT = {
    "l": "Left",
    "ctr": "Center",
    "r": "Right",
    "just": "Justify",
    "justLow": "JustifyLow",
    "dist": "Distributed",
    "thaiDist": "ThaiDistributed",
}

# a:bodyPr anchor (+ anchorCtr) -> PowerPoint.TextVerticalAlignment's exact
# string values (Top/Middle/Bottom/TopCentered/MiddleCentered/
# BottomCentered). The "Centered" variants correspond to anchorCtr="1" (the
# whole text block is also centered horizontally as a group) -- a separate
# OOXML attribute from paragraph-level algn, combined here since Office.js
# models both as one enum.
ANCHOR_TO_OFFICEJS_VERTICAL_ALIGNMENT = {"t": "Top", "ctr": "Middle", "b": "Bottom"}

# Raw OOXML <a:prstDash val="..."> -> Office.js ShapeLineDashStyle, by
# literal name correspondence -- NOT via python-pptx's MSO_LINE_DASH_STYLE,
# which was tried first and confirmed wrong: that enum's ROUND_DOT/
# SQUARE_DOT members are python-pptx's own reading of the legacy VBA
# msoLineDashStyle names, which round-trip to the *system* dash variants
# (raw "sysDot"/"sysDash" -- confirmed via MSO_LINE_DASH_STYLE.to_xml() and
# .from_xml(), both directions agree). Office.js's dashStyle enum has
# SystemDash/SystemDot/SystemDashDot as their own distinct, separately-named
# values alongside RoundDot/SquareDot, which only makes sense if Office.js's
# enum is keyed to the literal raw OOXML names (sys* -> System*) rather than
# the legacy VBA aliasing python-pptx uses -- confirmed as the actual bug
# after a real round-trip test showed a source line saved as "Square Dot"
# (raw prstDash="sysDash") pasting as something other than square-dotted
# once mapped through the old SQUARE_DOT->"SquareDot" table. Read directly
# off line._ln's raw XML instead of via line.dash_style, bypassing
# python-pptx's enum entirely so this table is the only translation layer.
# sysDashDotDot has no Office.js equivalent by name; DashDotDot is the only
# remaining unclaimed Office.js value, so it's the least-wrong fallback
# rather than dropping the dash style entirely.
PRST_DASH_TO_OFFICEJS = {
    "solid": "Solid",
    "dot": "RoundDot",
    "dash": "Dash",
    "lgDash": "LongDash",
    "dashDot": "DashDot",
    "lgDashDot": "LongDashDot",
    "lgDashDotDot": "LongDashDotDot",
    "sysDash": "SystemDash",
    "sysDot": "SystemDot",
    "sysDashDot": "SystemDashDot",
    "sysDashDotDot": "DashDotDot",
}

# <a:ln cmpd="..."> (compound/parallel-line style -- a different OOXML
# attribute from dash style) -> PowerPoint.ShapeLineStyle. Not exposed by
# python-pptx's public LineFormat API at all, read directly off line._ln
# below. "tri" (a thin-thick-thin triple line) is the best conceptual
# match for Office.js's "ThickBetweenThin" (a thick line sandwiched
# between two thin ones) -- the two enums don't share a common ancestor
# the way dash style's does, so this mapping is a best-effort semantic
# match, not a verified 1:1 correspondence.
CMPD_TO_OFFICEJS = {
    "sng": "Single",
    "dbl": "ThinThin",
    "thickThin": "ThickThin",
    "thinThick": "ThinThick",
    "tri": "ThickBetweenThin",
}

# MSO_UNDERLINE member name -> PowerPoint.ShapeFontUnderlineStyle. Built by
# hand rather than via to_xml() (unlike dash style above, the member names
# themselves already carry a direct "strip the _LINE suffix" correspondence
# to Office.js's names, so a hand-written table is simpler and no less
# accurate than round-tripping through the raw <a:u val="..."> value).
# WORDS (underline non-space runs only) and MIXED have no Office.js
# equivalent and are intentionally absent -- left unmapped rather than
# guessed, per this file's established "omit over guess" convention.
MSO_UNDERLINE_TO_OFFICEJS = {
    "NONE": "None",
    "SINGLE_LINE": "Single",
    "DOUBLE_LINE": "Double",
    "HEAVY_LINE": "Heavy",
    "DOTTED_LINE": "Dotted",
    "DOTTED_HEAVY_LINE": "DottedHeavy",
    "DASH_LINE": "Dash",
    "DASH_HEAVY_LINE": "DashHeavy",
    "DASH_LONG_LINE": "DashLong",
    "DASH_LONG_HEAVY_LINE": "DashLongHeavy",
    "DOT_DASH_LINE": "DotDash",
    "DOT_DASH_HEAVY_LINE": "DotDashHeavy",
    "DOT_DOT_DASH_LINE": "DotDotDash",
    "DOT_DOT_DASH_HEAVY_LINE": "DotDotDashHeavy",
    "WAVY_LINE": "Wavy",
    "WAVY_HEAVY_LINE": "WavyHeavy",
    "WAVY_DOUBLE_LINE": "WavyDouble",
}

# MSO_AUTO_SIZE member name -> PowerPoint.ShapeAutoSize. Note the two
# libraries name the "shape grows/shrinks" and "text shrinks" cases the
# opposite way round from each other (python-pptx's SHAPE_TO_FIT_TEXT --
# the shape resizes -- is Office.js's autoSizeShapeToFitText; python-pptx's
# TEXT_TO_FIT_SHAPE -- the text shrinks -- is autoSizeTextToFitShape) --
# confirmed by reading both enums' own descriptions, not assumed from the
# similar-looking names. MIXED has no single value.
MSO_AUTO_SIZE_TO_OFFICEJS = {
    "NONE": "AutoSizeNone",
    "SHAPE_TO_FIT_TEXT": "AutoSizeShapeToFitText",
    "TEXT_TO_FIT_SHAPE": "AutoSizeTextToFitShape",
}

# <a:buAutoNum type="..."> (a PowerPoint built-in auto-numbering scheme,
# e.g. "1. 2. 3." or "a. b. c.") -> PowerPoint.BulletFormat's style enum,
# which supports exactly this same fixed set of numbering schemes (unlike
# <a:buChar>, an arbitrary custom glyph, which BulletFormat cannot express
# at all -- that's still routed to 'file' mode unconditionally). Only the
# common Western/CJK/Thai schemes with an unambiguous Office.js
# counterpart are mapped; a handful of OOXML ST_TextAutonumberScheme
# values (hebrew2Minus and others) are deliberately left out rather than
# guessed -- an unmapped buAutoNum type routes the whole shape to 'file'
# mode (see classify_shape_tree), the same conservative fallback already
# used for buChar.
BU_AUTONUM_TO_OFFICEJS = {
    "alphaLcParenBoth": "AlphabetLowercaseParenthesesBoth",
    "alphaLcParenR": "AlphabetLowercaseParenthesisRight",
    "alphaLcPeriod": "AlphabetLowercasePeriod",
    "alphaUcParenBoth": "AlphabetUppercaseParenthesesBoth",
    "alphaUcParenR": "AlphabetUppercaseParenthesisRight",
    "alphaUcPeriod": "AlphabetUppercasePeriod",
    "arabicParenBoth": "ArabicNumeralParenthesesBoth",
    "arabicParenR": "ArabicNumeralParenthesisRight",
    "arabicPeriod": "ArabicNumeralPeriod",
    "arabicPlain": "ArabicNumeralPlain",
    "romanLcParenBoth": "RomanLowercaseParenthesesBoth",
    "romanLcParenR": "RomanLowercaseParenthesisRight",
    "romanLcPeriod": "RomanLowercasePeriod",
    "romanUcParenBoth": "RomanUppercaseParenthesesBoth",
    "romanUcParenR": "RomanUppercaseParenthesisRight",
    "romanUcPeriod": "RomanUppercasePeriod",
    "circleNumDbPlain": "CircleNumberDoubleBytePlain",
    "circleNumWdBlackPlain": "CircleNumberWideDoubleByteBlackPlain",
    "circleNumWdWhitePlain": "CircleNumberWideDoubleByteWhitePlain",
    "arabicDbPeriod": "ArabicDoubleBytePeriod",
    "arabicDbPlain": "ArabicDoubleBytePlain",
    "ea1ChsPeriod": "SimplifiedChinesePeriod",
    "ea1ChsPlain": "SimplifiedChinesePlain",
    "ea1ChtPeriod": "TraditionalChinesePeriod",
    "ea1ChtPlain": "TraditionalChinesePlain",
    "thaiAlphaParenBoth": "ThaiAlphabetParenthesesBoth",
    "thaiAlphaParenR": "ThaiAlphabetParenthesisRight",
    "thaiAlphaPeriod": "ThaiAlphabetPeriod",
    "thaiNumParenBoth": "ThaiNumeralParenthesesBoth",
    "thaiNumParenR": "ThaiNumeralParenthesisRight",
    "thaiNumPeriod": "ThaiNumeralPeriod",
}

# Fill types with no Office.js ShapeFill setter at all (only setSolidColor
# and setImage exist -- confirmed against @types/office-js) -- routed to
# 'file' mode the same as custom geometry, rather than silently reaching
# extract_fill and coming back None (today's actual bug: these shapes
# currently classify as 'reconstruct' and lose their fill entirely).
# PICTURE is deliberately included here too, even though setImage() could
# in principle carry it across -- extracting the underlying image blob
# needs reaching past python-pptx's public FillFormat API into relationship
# resolution with no way to verify the round-trip without a live
# PowerPoint test, so the safer, already-proven file-mode path is used
# instead, exactly like a real <p:pic> object already is.
#
# MSO_FILL_TYPE.BACKGROUND does NOT belong here, despite the name sounding
# like "a complex background fill" -- it's what python-pptx reports for a
# shape with an explicit <a:noFill/>, i.e. plain "no fill at all"
# (confirmed directly against python-pptx's own _NoFill.type source). That
# is fully reconstructable (Office.js: shape.fill.clear()) and is also the
# single most common fill state in real content -- PowerPoint writes
# <a:noFill/> explicitly on the overwhelming majority of plain text boxes.
# Including it here was a real bug: it silently routed nearly every plain
# text box (any shape with no fill at all) to 'file' mode, confirmed by
# reproducing against the actual Text.pptx boilerplate, where every single
# slide -- including plain single-text-box ones with nothing else going
# on -- classified as 'file' until this was removed.
NON_RECONSTRUCTABLE_FILL_TYPES = {
    MSO_FILL_TYPE.GRADIENT,
    MSO_FILL_TYPE.PATTERNED,
    MSO_FILL_TYPE.PICTURE,
    MSO_FILL_TYPE.TEXTURED,
}


def has_unreconstructable_fill(shape) -> bool:
    try:
        if shape.fill.type in NON_RECONSTRUCTABLE_FILL_TYPES:
            return True
    except Exception:
        pass
    # No explicit fill on the shape itself -- if it's coming from a
    # <p:style><a:fillRef> quick-style pointing at a gradient format-scheme
    # entry, that can't be flattened to one color either (see
    # has_unreconstructable_style_ref_fill, defined further down alongside
    # the rest of the style-ref resolution logic).
    return has_unreconstructable_style_ref_fill(shape)


def has_unmapped_auto_number(element) -> bool:
    """True if any paragraph under this shape uses a buAutoNum numbering
    scheme this app can't translate to an Office.js BulletFormat.style
    value (see BU_AUTONUM_TO_OFFICEJS) -- routes the whole shape to 'file'
    mode rather than silently dropping that paragraph's numbering, the
    same conservative fallback already used for buChar."""
    for bu in element.findall(".//" + qn("a:buAutoNum")):
        if bu.get("type") not in BU_AUTONUM_TO_OFFICEJS:
            return True
    return False


def classify_shape_tree(shape) -> str:
    """Returns 'file' if anything under this shape can't be reconstructed, else 'reconstruct'."""
    el = shape._element
    if (
        has_cust_geom(el)
        or has_picture(el)
        or has_other_graphic_frame(el)
        or has_custom_bullet_char(el)
        or has_unreconstructable_fill(shape)
        or has_unmapped_auto_number(el)
    ):
        return "file"
    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        return "reconstruct"
    prst_el = el.find(".//" + qn("a:prstGeom"))
    prst = prst_el.get("prst") if prst_el is not None else None
    if prst not in PRST_TO_GEOMETRIC_SHAPE_TYPE:
        return "file"
    return "reconstruct"


# ---------------------------------------------------------------------------
# Theme-color resolution. A color set via this app's own Fill/Line/Text
# color picker (src/features/fillLineColors.ts) -- or via PowerPoint's own
# theme-based color swatches, which is the common case for any shape not
# explicitly given a custom RGB color -- is stored as <a:schemeClr val="..">
# (e.g. "accent1"), not a literal RGB value. python-pptx's ColorFormat.rgb
# raises AttributeError on this color type ("no .rgb property on color type
# '_SchemeColor'") -- confirmed directly, and without handling it here every
# theme-colored fill/line/font color was silently dropped (caught by a
# blanket except and turned into None), which is why "recreate" mode was
# losing exactly the styles reported: colors are almost always theme-based
# in real content, explicit RGB is the unusual case.
#
# Resolving requires walking the real OOXML relationship chain: the shape's
# slide -> slide layout -> slide master -> theme part (for the actual
# <a:clrScheme> RGB values, keyed by raw slot names dk1/lt1/dk2/lt2/
# accent1-6/hlink/folHlink) and the slide master's own <p:clrMap> (which
# remaps the four semantic bg1/tx1/bg2/tx2 names PowerPoint's UI/API uses to
# whichever raw slot the current master actually points them at -- a "dark"
# master can map bg1 to dk1 instead of lt1. accent1-6/hlink/folHlink map to
# themselves, but only by convention; only the actual clrMap is Definitive.
# Verified against a real theme-colored test file before writing this: the
# resolved hex values matched the theme's own <a:clrScheme> exactly.
# ---------------------------------------------------------------------------


# Named preset colors (<a:prstClr val="...">) that a theme's <a:clrScheme>
# could in principle use in place of a literal RGB value -- rare, but a
# custom/converted template is exactly the kind of file likely to. Not the
# full CSS/ST_PresetColorVal list (150+ names), just enough common ones to
# avoid dropping a color entirely when this form is actually used; anything
# not in this small table still falls through to the theme-role-only
# fallback in resolve_color_spec rather than guessing.
PRST_CLR_TO_HEX = {
    "black": "000000", "white": "FFFFFF", "red": "FF0000", "green": "008000",
    "blue": "0000FF", "yellow": "FFFF00", "gray": "808080", "grey": "808080",
    "darkGray": "A9A9A9", "darkGrey": "A9A9A9", "lightGray": "D3D3D3", "lightGrey": "D3D3D3",
    "orange": "FFA500", "purple": "800080", "teal": "008080", "navy": "000080",
    "maroon": "800000", "olive": "808000", "silver": "C0C0C0", "aqua": "00FFFF",
    "fuchsia": "FF00FF", "lime": "00FF00",
}


def _hsl_clr_to_hex(el):
    """Converts an <a:hslClr hue="0-21600000" sat="0-100000" lum="0-100000">
    element to a #RRGGBB hex string -- one of three ways OOXML can express a
    theme color (alongside srgbClr and sysClr, already handled), seen in
    templates authored or converted via tools that store colors as HSL
    rather than RGB. hue is in 60,000ths of a degree; sat/lum are in
    1,000ths of a percent -- both confirmed against the OOXML DrawingML
    schema (CT_HslColor), not assumed from the attribute names alone."""
    try:
        hue = (int(el.get("hue")) / 60000.0) / 360.0
        sat = int(el.get("sat")) / 100000.0
        lum = int(el.get("lum")) / 100000.0
    except (TypeError, ValueError):
        return None
    r, g, b = colorsys.hls_to_rgb(hue, lum, sat)
    return f"{round(r * 255):02X}{round(g * 255):02X}{round(b * 255):02X}"


def _resolve_scheme_color_element(child):
    """A <a:clrScheme> slot's color can be expressed as srgbClr, sysClr,
    hslClr, or prstClr -- returns the raw hex (no '#') for whichever form
    is actually present, or None if it's some other/unrecognized form
    (e.g. scrgbClr, a scientific-RGB form no real template has been seen
    using for a theme color)."""
    srgb = child.find(qn("a:srgbClr"))
    if srgb is not None:
        return srgb.get("val")
    sys_clr = child.find(qn("a:sysClr"))
    if sys_clr is not None:
        return sys_clr.get("lastClr")
    hsl_clr = child.find(qn("a:hslClr"))
    if hsl_clr is not None:
        return _hsl_clr_to_hex(hsl_clr)
    prst_clr = child.find(qn("a:prstClr"))
    if prst_clr is not None:
        return PRST_CLR_TO_HEX.get(prst_clr.get("val"))
    return None


def _load_master_theme(master):
    """Returns (clr_map, clr_scheme) for a slide master, both plain dicts of
    XML slot name -> value. Cached per (per-request) theme_cache dict, since
    a whole category's worth of slides typically share one master and
    re-walking the relationship chain/re-parsing the theme XML per shape
    would be wasteful."""
    clr_map_el = master.element.find(qn("p:clrMap"))
    clr_map = dict(clr_map_el.attrib) if clr_map_el is not None else {}

    theme_part = None
    for rel in master.part.rels.values():
        if rel.reltype.endswith("/theme"):
            theme_part = rel.target_part
            break

    scheme = {}
    if theme_part is not None:
        theme_el = etree.fromstring(theme_part.blob)
        clr_scheme_el = theme_el.find(".//" + qn("a:clrScheme"))
        if clr_scheme_el is not None:
            for child in clr_scheme_el:
                tag = etree.QName(child).localname
                hexval = _resolve_scheme_color_element(child)
                if hexval:
                    scheme[tag] = hexval

    return clr_map, scheme


# Raw OOXML schemeClr value (pre-clrMap -- MSO_THEME_COLOR.to_xml() returns
# exactly this, the semantic role name PowerPoint's own UI/API uses, not the
# raw theme-XML slot it happens to map to on this particular master) ->
# Office.js's ThemeColorScheme.getThemeColor() role names, confirmed
# against src/features/fillLineColors.ts's own THEME_COLOR_ROLES list.
# hlink/folHlink are deliberately absent: that file's getThemeColors() never
# exposes them either, so there's no live-lookup path for a hyperlink color
# at this app's PowerPointApi floor -- those colors fall back to the static
# hex snapshot below and nothing else.
#
# MSO_THEME_COLOR.to_xml() can also return the four RAW scheme slot names
# (dk1/lt1/dk2/lt2) instead of the semantic ones, when a shape's
# <a:schemeClr val="..."> happens to reference a slot directly rather than
# through clrMap (confirmed directly: python-pptx round-trips whichever
# form was actually written -- DARK_1/LIGHT_1/etc. produce "dk1"/"lt1"/...,
# while TEXT_1/BACKGROUND_1/etc. produce "tx1"/"bg1"/...). Mapped here to
# the same Office.js role their semantic counterpart would typically
# resolve to, since tx1->dk1 and bg1->lt1 is clrMap's standard/default
# configuration -- an approximation for the (uncommon) direct-slot case,
# not a guess for the common semantic case, which is exact.
XML_NAME_TO_OFFICEJS_ROLE = {
    "bg1": "Light1",
    "tx1": "Dark1",
    "bg2": "Light2",
    "tx2": "Dark2",
    "lt1": "Light1",
    "dk1": "Dark1",
    "lt2": "Light2",
    "dk2": "Dark2",
    "accent1": "Accent1",
    "accent2": "Accent2",
    "accent3": "Accent3",
    "accent4": "Accent4",
    "accent5": "Accent5",
    "accent6": "Accent6",
}


def resolve_theme_color(shape, mso_theme_color, theme_cache):
    """Returns (hex, officejs_role). hex is a static snapshot of the SOURCE
    presentation's theme at extraction time -- kept as a fallback. When
    officejs_role is non-None, the insert-time client can instead resolve
    the color live against the TARGET presentation's own theme (see
    src/features/libraryInsert.ts's two-phase theme-role resolve), so a
    shape colored e.g. "Accent1" in the source lands as whatever Accent1
    actually is in the presentation it's inserted into, rather than a
    baked-in hex that's wrong (commonly white-on-white) the moment the two
    presentations' themes differ."""
    try:
        xml_name = MSO_THEME_COLOR.to_xml(mso_theme_color)
    except Exception:
        return None, None
    officejs_role = XML_NAME_TO_OFFICEJS_ROLE.get(xml_name)
    try:
        master = shape.part.slide_layout.slide_master
    except Exception:
        return None, officejs_role
    cache_key = id(master.part)
    if cache_key not in theme_cache:
        theme_cache[cache_key] = _load_master_theme(master)
    clr_map, scheme = theme_cache[cache_key]
    slot = clr_map.get(xml_name, xml_name)
    hexval = scheme.get(slot)
    return (f"#{hexval}" if hexval else None), officejs_role


def resolve_color_spec(color_format, shape, theme_cache):
    """Resolves a python-pptx ColorFormat to {"hex", "themeRole"}: hex a
    real #RRGGBB string, themeRole a live Office.js theme-color role name
    when this color is theme-derived and mappable, else None. Returns None
    entirely if no color is set, or for the rare color types (preset/
    scRGB) this doesn't handle -- safer to omit than to guess.

    Bug fixed here: previously required BOTH hex and role to resolve
    before keeping the color at all ("if hexval else None"), which threw
    away a perfectly good theme role -- and with it, the shape's entire
    fill/line/font color -- whenever the static hex snapshot alone failed
    to compute (e.g. a custom-named brand color like "Dark Teal" defined
    via <a:hslClr> rather than <a:srgbClr>/<a:sysClr>, before those were
    handled in _load_master_theme). A resolvable role is enough on its own:
    the insert-time client resolves it live against the target
    presentation's actual theme (see libraryInsert.ts's resolveThemeRoles)
    and never even looks at hex unless live resolution is unavailable. A
    neutral gray placeholder hex is used only for that fallback case,
    instead of dropping the color (and silently rendering as white-on-
    white, or here, no fill at all) whenever the snapshot alone fails."""
    try:
        color_type = color_format.type
    except Exception:
        return None
    if color_type == MSO_COLOR_TYPE.RGB:
        return {"hex": f"#{color_format.rgb}", "themeRole": None}
    if color_type == MSO_COLOR_TYPE.SCHEME:
        hexval, role = resolve_theme_color(shape, color_format.theme_color, theme_cache)
        if not hexval and not role:
            return None
        return {"hex": hexval or "#808080", "themeRole": role}
    return None


# ---------------------------------------------------------------------------
# "Shape Style" quick-formatting (<p:style><a:lnRef>/<a:fillRef>/<a:fontRef>).
# Confirmed as the actual root cause of a real "theme colour pastes as no
# colour" report: when a shape has no EXPLICIT <a:solidFill>/<a:ln>/run-
# level <a:solidFill> of its own, PowerPoint still renders it using
# whichever color the shape's <p:style> block references -- this is how a
# freshly-drawn shape gets its default fill/line/font color, and it's a
# completely separate mechanism from the direct <a:schemeClr> references
# resolve_theme_color already handles. shape.fill.type/shape.line.fill.type/
# run.font.color.type are all None in this case (python-pptx only reads
# spPr/rPr directly, never <p:style>), which is why extraction was silently
# reporting "no fill" / "no line" / "no run color" for shapes that actually
# render with a real color in PowerPoint -- confirmed directly against a
# real user-reported file (a freshly-drawn oval with fillRef idx="1"
# accent1 and no spPr fill at all).
# ---------------------------------------------------------------------------


def _find_style_ref(shape, ref_tag):
    """Returns the <p:style><a:{ref_tag}> element for this shape, or None
    if the shape has no <p:style> block at all (e.g. a plain text box,
    which never gets one) or no ref of that kind."""
    style_el = shape._element.find(qn("p:style"))
    if style_el is None:
        return None
    return style_el.find(qn("a:" + ref_tag))


def _apply_color_modifiers(hexval, scheme_clr_el):
    """Applies the common OOXML color transforms (shade/tint/lumMod/
    lumOff) that can appear as direct children of a <p:style> ref's
    <a:schemeClr>, adjusting an already-resolved base hex via simple HSL
    luminance math. Anything more exotic (satMod, hueMod, etc.) is left
    unapplied -- the base color is still far better than nothing, and
    those are rare on a shape-level style ref in practice (they're far
    more common inside the theme's own gradient stop definitions, which
    this app already refuses to flatten -- see
    _format_scheme_entry_is_gradient below)."""
    if scheme_clr_el is None or hexval is None:
        return hexval
    try:
        r = int(hexval[1:3], 16) / 255.0
        g = int(hexval[3:5], 16) / 255.0
        b = int(hexval[5:7], 16) / 255.0
    except (ValueError, IndexError):
        return hexval
    h, l, s = colorsys.rgb_to_hls(r, g, b)
    for tag, apply in (
        ("a:shade", lambda v, l: l * (v / 100000.0)),
        ("a:tint", lambda v, l: l * (v / 100000.0) + (1.0 - v / 100000.0)),
        ("a:lumMod", lambda v, l: l * (v / 100000.0)),
        ("a:lumOff", lambda v, l: l + v / 100000.0),
    ):
        el = scheme_clr_el.find(qn(tag))
        if el is not None:
            try:
                l = apply(int(el.get("val")), l)
            except (TypeError, ValueError):
                pass
    l = max(0.0, min(1.0, l))
    r2, g2, b2 = colorsys.hls_to_rgb(h, l, s)
    return f"#{round(r2 * 255):02X}{round(g2 * 255):02X}{round(b2 * 255):02X}"


def _format_scheme_entry(master, list_tag, idx):
    """Returns the theme's <a:fmtScheme><a:{list_tag}> entry at 1-based
    idx (an lxml element), or None if idx is 0/missing/out of range. A
    style ref's idx indexes into this 3-entry list (fillStyleLst/
    lnStyleLst), which is where the format scheme -- not the shape itself
    -- defines whether that "look" is a flat solid or a gradient, and
    (for lines) the actual weight."""
    if not idx or idx == "0":
        return None
    try:
        theme_part = None
        for rel in master.part.rels.values():
            if rel.reltype.endswith("/theme"):
                theme_part = rel.target_part
                break
        if theme_part is None:
            return None
        theme_el = etree.fromstring(theme_part.blob)
        list_el = theme_el.find(".//" + qn("a:fmtScheme") + "/" + qn("a:" + list_tag))
        if list_el is None:
            return None
        entries = list(list_el)
        return entries[int(idx) - 1]
    except Exception:
        return None


def _format_scheme_entry_is_gradient(master, list_tag, idx):
    entry = _format_scheme_entry(master, list_tag, idx)
    return entry is not None and entry.tag == qn("a:gradFill")


def resolve_style_ref_spec(shape, ref_tag, theme_cache):
    """Resolves a <p:style> quick-style reference to a {"hex", "themeRole"}
    spec, the same shape resolve_color_spec returns -- used as the
    fallback for fill/line/font color whenever the shape/run has no
    EXPLICIT color of its own (see extract_fill/extract_line/
    extract_text_spec). Returns "gradient" as a sentinel when the
    referenced fillStyleLst entry is a gradient (fontRef has no such
    concept and never returns this), so the caller can route the whole
    shape to 'file' mode instead of flattening a gradient to one wrong
    flat color -- the same reason a shape's own direct gradient fill
    already routes to 'file'."""
    ref_el = _find_style_ref(shape, ref_tag)
    if ref_el is None:
        return None
    idx = ref_el.get("idx")
    if not idx or idx == "0":
        return None
    scheme_clr_el = ref_el.find(qn("a:schemeClr"))
    if scheme_clr_el is None:
        return None
    xml_name = scheme_clr_el.get("val")
    if not xml_name:
        return None
    try:
        master = shape.part.slide_layout.slide_master
    except Exception:
        return None
    list_tag = {"fillRef": "fillStyleLst", "lnRef": "lnStyleLst"}.get(ref_tag)
    if list_tag and _format_scheme_entry_is_gradient(master, list_tag, idx):
        return "gradient"
    officejs_role = XML_NAME_TO_OFFICEJS_ROLE.get(xml_name)
    cache_key = id(master.part)
    if cache_key not in theme_cache:
        theme_cache[cache_key] = _load_master_theme(master)
    clr_map, scheme = theme_cache[cache_key]
    slot = clr_map.get(xml_name, xml_name)
    hexval = scheme.get(slot)
    hexval = f"#{hexval}" if hexval else None
    hexval = _apply_color_modifiers(hexval, scheme_clr_el)
    if not hexval and not officejs_role:
        return None
    return {"hex": hexval or "#808080", "themeRole": officejs_role}


def style_ref_line_width_pt(shape):
    """The line weight the format scheme's lnStyleLst[idx-1] entry
    specifies for this shape's <p:style><a:lnRef> -- used only when the
    shape has no explicit <a:ln> of its own (see extract_line), since
    there's no other source for a sensible width in that case."""
    ref_el = _find_style_ref(shape, "lnRef")
    if ref_el is None:
        return None
    idx = ref_el.get("idx")
    if not idx or idx == "0":
        return None
    try:
        master = shape.part.slide_layout.slide_master
    except Exception:
        return None
    entry = _format_scheme_entry(master, "lnStyleLst", idx)
    if entry is None:
        return None
    w = entry.get("w")
    return emu_to_pt(int(w)) if w else None


def has_unreconstructable_style_ref_fill(shape) -> bool:
    """True when a shape with no explicit fill relies on a <p:style>
    fillRef pointing at a gradient format-scheme entry -- classify_shape_
    tree's counterpart to has_unreconstructable_fill, for the style-ref
    case that function can't see (shape.fill.type is None there, not one
    of the gradient/pattern/etc. enum members)."""
    try:
        if shape.fill.type is not None:
            return False
        ref_el = _find_style_ref(shape, "fillRef")
        if ref_el is None:
            return False
        master = shape.part.slide_layout.slide_master
        return _format_scheme_entry_is_gradient(master, "fillStyleLst", ref_el.get("idx"))
    except Exception:
        return False


def extract_alpha(color_format):
    """Returns transparency as a 0.0 (opaque) - 1.0 (fully clear) fraction,
    matching Office.js's ShapeFill.transparency/ShapeLineFormat.transparency
    convention -- read from OOXML's <a:alpha val="N"/> (N = thousandths of
    a percent OPACITY, 100000 = fully opaque), which isn't exposed by
    python-pptx's public ColorFormat API at all, so read directly off the
    color's own raw XML element (color_format._color._xClr -- the same
    private-internals pattern already used elsewhere in this file, e.g.
    shape._element, line._ln). Returns 0.0 (opaque, i.e. no-op when applied)
    for any color type/internal shape this can't reach, rather than raise."""
    try:
        color_el = color_format._color._xClr
    except Exception:
        return 0.0
    alpha_el = color_el.find(qn("a:alpha"))
    if alpha_el is None:
        return 0.0
    try:
        return max(0.0, min(1.0, 1.0 - int(alpha_el.get("val")) / 100000.0))
    except (TypeError, ValueError):
        return 0.0


def extract_fill(shape, theme_cache):
    try:
        fill = shape.fill
        if fill.type is None:
            # No explicit fill on the shape itself -- fall back to the
            # <p:style><a:fillRef> quick-style color, if any (see the
            # style-ref resolution block above for why this is common,
            # not an edge case).
            style_color = resolve_style_ref_spec(shape, "fillRef", theme_cache)
            if style_color and style_color != "gradient":
                return {"type": "solid", "color": style_color, "transparency": 0.0}
            return None
        if str(fill.type) == "MSO_FILL_TYPE.SOLID (1)" or fill.type == 1:
            color = resolve_color_spec(fill.fore_color, shape, theme_cache)
            if not color:
                return None
            return {"type": "solid", "color": color, "transparency": extract_alpha(fill.fore_color)}
        return None
    except Exception:
        return None


def extract_line(shape, theme_cache):
    try:
        line = shape.line
        if line.fill.type is None:
            # No explicit line on the shape itself -- fall back to the
            # <p:style><a:lnRef> quick-style color/width, if any.
            style_color = resolve_style_ref_spec(shape, "lnRef", theme_cache)
            if style_color and style_color != "gradient":
                return {
                    "color": style_color,
                    "widthPt": style_ref_line_width_pt(shape),
                    "transparency": 0.0,
                    "dashStyle": None,
                    "compoundStyle": None,
                }
            return None
        width_pt = emu_to_pt(line.width) if line.width else None
        color = resolve_color_spec(line.color, shape, theme_cache)
        if not color:
            return None
        dash_style = None
        compound_style = None
        try:
            prst_dash_el = line._ln.find(qn("a:prstDash"))
            if prst_dash_el is not None:
                dash_style = PRST_DASH_TO_OFFICEJS.get(prst_dash_el.get("val"))
        except Exception:
            dash_style = None
        try:
            compound_style = CMPD_TO_OFFICEJS.get(line._ln.get("cmpd"))
        except Exception:
            compound_style = None
        return {
            "color": color,
            "widthPt": width_pt,
            "transparency": extract_alpha(line.color),
            "dashStyle": dash_style,
            "compoundStyle": compound_style,
        }
    except Exception:
        return None


def extract_bullet(p):
    """Returns {"style": ...} for a paragraph using one of Office.js's
    built-in auto-numbering schemes (<a:buAutoNum>), else None. A shape
    with a buAutoNum whose type isn't in BU_AUTONUM_TO_OFFICEJS never
    reaches here at all -- classify_shape_tree already routed it to 'file'
    mode (has_unmapped_auto_number), the same as it already does for
    <a:buChar>'s arbitrary custom glyphs, which BulletFormat has no way to
    express regardless."""
    pPr = p._p.find(qn("a:pPr"))
    if pPr is None:
        return None
    bu = pPr.find(qn("a:buAutoNum"))
    if bu is None:
        return None
    style = BU_AUTONUM_TO_OFFICEJS.get(bu.get("type"))
    return {"style": style} if style else None


def extract_text_spec(shape, theme_cache):
    if not shape.has_text_frame:
        return None
    paragraphs = []
    for p in shape.text_frame.paragraphs:
        runs = []
        for r in p.runs:
            underline = None
            try:
                if r.font.underline is not None and r.font.underline is not True and r.font.underline is not False:
                    underline = MSO_UNDERLINE_TO_OFFICEJS.get(r.font.underline.name)
                elif r.font.underline is True:
                    underline = "Single"
                elif r.font.underline is False:
                    underline = "None"
            except Exception:
                underline = None
            run_color = resolve_color_spec(r.font.color, shape, theme_cache) if r.font.color else None
            if run_color is None:
                # No explicit color on this run -- fall back to the
                # shape-level <p:style><a:fontRef> quick-style color (e.g.
                # a shape whose text was never given its own color
                # override, relying on the default that comes with the
                # shape's style).
                style_color = resolve_style_ref_spec(shape, "fontRef", theme_cache)
                if style_color and style_color != "gradient":
                    run_color = style_color
            runs.append(
                {
                    "text": r.text,
                    "bold": r.font.bold,
                    "italic": r.font.italic,
                    "underline": underline,
                    "size": r.font.size.pt if r.font.size else None,
                    "fontName": r.font.name,
                    "color": run_color,
                }
            )
        align = None
        try:
            if p.alignment is not None:
                align = ALGN_TO_OFFICEJS_ALIGNMENT.get(PP_ALIGN.to_xml(p.alignment))
        except Exception:
            align = None
        paragraphs.append({"level": p.level, "bullet": extract_bullet(p), "align": align, "runs": runs})
    return paragraphs


def extract_vertical_alignment(shape):
    """Maps a:bodyPr anchor/anchorCtr to Office.js's TextVerticalAlignment
    string. Returns None (not a guessed default) when the shape has no text
    frame or no anchor attribute at all -- the caller only sets
    TextFrame.verticalAlignment when this is non-null, leaving PowerPoint's
    own default (Top) untouched otherwise."""
    if not shape.has_text_frame:
        return None
    body_pr = shape._element.find(".//" + qn("a:bodyPr"))
    if body_pr is None:
        return None
    base = ANCHOR_TO_OFFICEJS_VERTICAL_ALIGNMENT.get(body_pr.get("anchor"))
    if base is None:
        return None
    return f"{base}Centered" if body_pr.get("anchorCtr") == "1" else base


def extract_text_frame_props(shape):
    """wordWrap/autoSize/margins -- all real python-pptx TextFrame
    properties (word_wrap, auto_size, margin_left/right/top/bottom), just
    never threaded through to the reconstruct spec until now. Returns a
    dict of Nones (not applied at insert time -- see buildShape) when the
    shape has no text frame, or for whichever individual fields aren't
    set/mappable, rather than guessing a default."""
    if not shape.has_text_frame:
        return {"wordWrap": None, "autoSize": None, "marginLeft": None, "marginRight": None, "marginTop": None, "marginBottom": None}
    tf = shape.text_frame
    auto_size = None
    try:
        if tf.auto_size is not None:
            auto_size = MSO_AUTO_SIZE_TO_OFFICEJS.get(tf.auto_size.name)
    except Exception:
        auto_size = None
    return {
        "wordWrap": tf.word_wrap,
        "autoSize": auto_size,
        "marginLeft": emu_to_pt(tf.margin_left),
        "marginRight": emu_to_pt(tf.margin_right),
        "marginTop": emu_to_pt(tf.margin_top),
        "marginBottom": emu_to_pt(tf.margin_bottom),
    }


def extract_adjustments(shape):
    """Adjustment-handle values (e.g. a rounded rectangle's corner radius,
    a chevron's arrow point, a callout's tail position -- <a:avLst><a:gd
    name="adj" fmla="val N"/></a:avLst>). Only autoshapes with adjustment
    handles expose shape.adjustments at all (python-pptx raises/lacks the
    attribute for a plain textbox or a shape with no adjustable geometry) --
    confirmed empirically that python-pptx's normalized fraction (e.g.
    0.16667 for a default rounded rectangle) and Office.js's
    Shape.adjustments.get/set use the identical 0.0-1.0 convention, so no
    unit conversion is needed, unlike almost everything else in this file."""
    try:
        values = list(shape.adjustments)
    except Exception:
        return None
    return values if values else None


def extract_reconstruct_spec(shape, theme_cache):
    is_text_box = shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    text_frame_props = extract_text_frame_props(shape)
    spec = {
        "kind": "textBox" if is_text_box else "geometricShape",
        "left": emu_to_pt(shape.left),
        "top": emu_to_pt(shape.top),
        "width": emu_to_pt(shape.width),
        "height": emu_to_pt(shape.height),
        "rotation": shape.rotation,
        "adjustments": extract_adjustments(shape),
        "fill": extract_fill(shape, theme_cache),
        "line": extract_line(shape, theme_cache),
        "verticalAlignment": extract_vertical_alignment(shape),
        "wordWrap": text_frame_props["wordWrap"],
        "autoSize": text_frame_props["autoSize"],
        "marginLeft": text_frame_props["marginLeft"],
        "marginRight": text_frame_props["marginRight"],
        "marginTop": text_frame_props["marginTop"],
        "marginBottom": text_frame_props["marginBottom"],
        "paragraphs": extract_text_spec(shape, theme_cache),
    }
    if not is_text_box:
        prst_el = shape._element.find(".//" + qn("a:prstGeom"))
        prst = prst_el.get("prst") if prst_el is not None else None
        spec["presetGeometry"] = PRST_TO_GEOMETRIC_SHAPE_TYPE[prst]
    return spec


@app.post("/convert")
def convert():
    upload = request.files.get("file")
    if upload is None or upload.filename == "":
        raise ConversionError("No file uploaded (expected multipart field 'file').")
    if not upload.filename.lower().endswith((".pptx", ".potx")):
        raise ConversionError("File must be a .pptx or .potx.")

    work_dir = tempfile.mkdtemp(prefix="render-")
    try:
        source_path = os.path.join(work_dir, "source.pptx")
        upload.save(source_path)

        try:
            prs = Presentation(source_path)
            slide_count = len(prs.slides)
        except Exception as exc:
            raise ConversionError(f"Could not read the file as a presentation: {exc}")
        if slide_count == 0:
            raise ConversionError("Presentation has no slides.")

        # Isolated LibreOffice profile per request -- a profile dir shared
        # across concurrent requests can lock/conflict.
        profile_dir = os.path.join(work_dir, "loprofile")
        run([
            "soffice", "--headless", "--invisible", "--nocrashreport",
            "--nodefault", "--nofirststartwizard", "--nolockcheck", "--nologo",
            "--norestore", f"-env:UserInstallation=file://{profile_dir}",
            "--convert-to", "pdf", "--outdir", work_dir, source_path,
        ])
        pdf_path = os.path.join(work_dir, "source.pdf")
        if not os.path.exists(pdf_path):
            raise ConversionError("LibreOffice did not produce a PDF.", status=422)

        run(["pdftoppm", "-png", "-r", str(THUMBNAIL_DPI), pdf_path, os.path.join(work_dir, "slide")])

        # pdftoppm zero-pads its page-number suffix based on total page
        # count (slide-1.png for <10 pages, slide-01.png for 10-99, etc.),
        # so a sorted glob -- not a guessed filename -- is the only correct
        # way to recover slide order for an arbitrary-sized deck.
        produced_pngs = sorted(
            name for name in os.listdir(work_dir) if re.fullmatch(r"slide-\d+\.png", name)
        )
        if len(produced_pngs) != slide_count:
            raise ConversionError(
                f"Expected {slide_count} rendered slide(s), got {len(produced_pngs)}.",
                status=422,
            )

        # One per request (not module-level) so a slide master's cached
        # theme never leaks across unrelated uploads -- see
        # _load_master_theme's docstring for why caching happens at all.
        theme_cache = {}

        manifest = {"slides": []}
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for index, raw_png_name in enumerate(produced_pngs):
                page_num = index + 1
                raw_png = os.path.join(work_dir, raw_png_name)
                trimmed_png = os.path.join(work_dir, f"trimmed-{page_num}.png")
                run([
                    "convert", raw_png, "-trim", "+repage",
                    "-bordercolor", "white", "-border", str(THUMBNAIL_MARGIN_PX),
                    trimmed_png,
                ])
                thumb_name = f"slide-{page_num}.png"
                zf.write(trimmed_png, thumb_name)

                slide = prs.slides[index]
                real_shapes = [
                    s for s in slide.shapes if not is_think_cell_placeholder(s) and not is_unused_layout_placeholder(s)
                ]

                entry = {
                    "index": index,
                    "title": extract_title_hint(slide),
                    "thumbnail": thumb_name,
                }

                if real_shapes and all(classify_shape_tree(s) == "reconstruct" for s in real_shapes):
                    if len(real_shapes) == 1:
                        spec = extract_reconstruct_spec(real_shapes[0], theme_cache)
                    else:
                        spec = {
                            "kind": "group",
                            "shapes": [extract_reconstruct_spec(s, theme_cache) for s in real_shapes],
                        }
                    entry["insertMode"] = "reconstruct"
                    entry["reconstructSpec"] = spec
                    entry["pptx"] = None
                else:
                    slide_pptx = os.path.join(work_dir, f"slide-{page_num}.pptx")
                    slice_single_slide(source_path, index, slide_pptx)
                    pptx_name = f"slide-{page_num}.pptx"
                    zf.write(slide_pptx, pptx_name)
                    entry["insertMode"] = "file"
                    entry["reconstructSpec"] = None
                    entry["pptx"] = pptx_name

                manifest["slides"].append(entry)
            zf.writestr("manifest.json", json.dumps(manifest))

        zip_buffer.seek(0)
        return send_file(zip_buffer, mimetype="application/zip", download_name="converted.zip")
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=8090)
