#!/usr/bin/env python3
"""
One-off content-prep script for Flags.pptx - structurally different from
every other category, so it isn't handled by the general
slice-catalog-source.py pipeline. Flags.pptx bundles many countries onto
one slide (e.g. "Africa A-G" = 15 country-name labels + 15 flag pictures
laid out in a 5-column grid), rather than one slide = one catalog item.
This script pairs each label with its flag picture and produces one
catalog item per country instead of one per slide.

Pairing (confirmed by inspecting real slide coordinates): each flag
picture sits directly above its matching country-name label, in the same
column. For each label, the nearest picture above it (by vertical
distance) whose horizontal center is within PAIR_TOLERANCE_PT is its
match; a picture can only be used once. A label with no match within
tolerance is skipped with a warning rather than crashing - the source
deck has at least one real data-quality issue (a duplicated, misspelled
"Gutemala" label with no picture of its own).

Matching happens per-container (the slide's own top-level shapes, and
separately each GROUP shape's own children, recursing arbitrarily deep)
since a group's child coordinates are only meaningful relative to that
same group, not the slide.

Each matched pair becomes its own single-picture single-slide .pptx
(960x540pt, matching every other category, not a truncated copy of the
busy multi-shape source slide) with the picture centered at its original
size, always 'file' mode (a picture can never reconstruct). Thumbnails
reuse the same qlmanage approach as slice-catalog-source.py, duplicated
here rather than imported - that script's hyphenated filename can't be
imported as a plain Python module, and this is a one-off job, not a
reusable pipeline step (matching the scripts/append-symbol-chars.js
precedent).

Usage:
    python3 scripts/slice-flags-source.py <Flags.pptx> [output-dir]
"""
import io
import json
import subprocess
import sys
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Pt

THINK_CELL_MARKER = "think-cell"
TARGET_SLIDE_WIDTH_PT = 960.0
TARGET_SLIDE_HEIGHT_PT = 540.0
RENDER_SIZE = 1600
CONTENT_PADDING_FRAC = 0.15
CONTENT_PADDING_MIN_PT = 8
PAIR_TOLERANCE_PT = 15.0


def is_think_cell_placeholder(shape) -> bool:
    return THINK_CELL_MARKER in shape._element.xml


def emu_to_pt(value) -> float:
    return round(Emu(value).pt, 2) if value is not None else None


def iter_pairing_containers(shapes):
    """
    Yields each "container" of shapes that should be paired independently:
    the given shapes themselves, plus (recursively) each GROUP's own
    children - a group's child coordinates are only meaningful relative
    to that same group, not the slide, so pairing must never mix shapes
    from two different containers.
    """
    yield [s for s in shapes if not is_think_cell_placeholder(s)]
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_pairing_containers(shape.shapes)


def pair_labels_and_pictures(container):
    """Returns a list of (label, picture) pairs for one container, per the module's pairing rule."""
    labels = [s for s in container if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.has_text_frame and s.text_frame.text.strip()]
    pictures = [s for s in container if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    used_picture_ids = set()
    pairs = []
    for label in labels:
        label_center_x = emu_to_pt(label.left) + emu_to_pt(label.width) / 2
        label_top = emu_to_pt(label.top)
        best, best_gap = None, None
        for pic in pictures:
            if id(pic) in used_picture_ids:
                continue
            pic_center_x = emu_to_pt(pic.left) + emu_to_pt(pic.width) / 2
            pic_top = emu_to_pt(pic.top)
            if abs(pic_center_x - label_center_x) > PAIR_TOLERANCE_PT:
                continue
            if pic_top >= label_top:
                continue
            gap = label_top - pic_top
            if best is None or gap < best_gap:
                best, best_gap = pic, gap
        if best is None:
            print(f'  warning: no matching picture found for label "{label.text_frame.text.strip()}", skipping', file=sys.stderr)
            continue
        used_picture_ids.add(id(best))
        pairs.append((label, best))
    return pairs


def build_flag_slide(picture, dest_path: Path) -> None:
    """Builds a brand-new 960x540pt single-slide presentation with just this one flag picture, centered at its original size."""
    prs = Presentation()
    prs.slide_width = Emu(round(TARGET_SLIDE_WIDTH_PT * 12700))
    prs.slide_height = Emu(round(TARGET_SLIDE_HEIGHT_PT * 12700))
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # "Blank"

    width_pt, height_pt = emu_to_pt(picture.width), emu_to_pt(picture.height)
    left_pt = (TARGET_SLIDE_WIDTH_PT - width_pt) / 2
    top_pt = (TARGET_SLIDE_HEIGHT_PT - height_pt) / 2

    image_stream = io.BytesIO(picture.image.blob)
    slide.shapes.add_picture(image_stream, Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(height_pt))

    dest_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest_path))


def generate_thumbnail(slide_pptx_path: Path, dest_path: Path, content_bbox_pt) -> bool:
    """Same approach as slice-catalog-source.py's generate_thumbnail (duplicated - see module docstring)."""
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            ["qlmanage", "-t", "-s", str(RENDER_SIZE), "-o", tmp, str(slide_pptx_path)],
            capture_output=True,
            check=False,
        )
        produced = Path(tmp) / f"{slide_pptx_path.name}.png"
        if not produced.exists():
            print(f"  warning: qlmanage produced no thumbnail for {slide_pptx_path.name}", file=sys.stderr)
            return False

        dest_path.parent.mkdir(parents=True, exist_ok=True)
        img = Image.open(produced)
        scale_x = img.width / TARGET_SLIDE_WIDTH_PT
        scale_y = img.height / TARGET_SLIDE_HEIGHT_PT
        left, top, right, bottom = content_bbox_pt
        box_px = (round(left * scale_x), round(top * scale_y), round(right * scale_x), round(bottom * scale_y))
        img.crop(box_px).save(dest_path)
        return True


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    source_path = Path(sys.argv[1])
    output_root = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("data/catalog")
    category = "flags"
    category_dir = output_root / category
    thumbnails_dir = output_root / "thumbnails" / category

    prs = Presentation(str(source_path))
    items = []
    item_number = 0

    for slide_index, slide in enumerate(prs.slides):
        real_shapes = [s for s in slide.shapes if not is_think_cell_placeholder(s)]
        pairs = []
        for container in iter_pairing_containers(real_shapes):
            pairs.extend(pair_labels_and_pictures(container))

        if not pairs:
            print(f"slide {slide_index + 1}: no label/picture pairs found, skipping", file=sys.stderr)
            continue

        for label, picture in pairs:
            item_number += 1
            title = label.text_frame.text.strip()
            slide_filename = f"{category}-{item_number:03d}.pptx"
            slide_pptx_path = category_dir / slide_filename
            build_flag_slide(picture, slide_pptx_path)

            width_pt, height_pt = emu_to_pt(picture.width), emu_to_pt(picture.height)
            left_pt = (TARGET_SLIDE_WIDTH_PT - width_pt) / 2
            top_pt = (TARGET_SLIDE_HEIGHT_PT - height_pt) / 2
            pad = max(CONTENT_PADDING_FRAC * max(width_pt, height_pt), CONTENT_PADDING_MIN_PT)
            content_bbox_pt = (
                max(0.0, left_pt - pad),
                max(0.0, top_pt - pad),
                min(TARGET_SLIDE_WIDTH_PT, left_pt + width_pt + pad),
                min(TARGET_SLIDE_HEIGHT_PT, top_pt + height_pt + pad),
            )
            thumbnail_filename = f"{category}-{item_number:03d}.png"
            has_thumbnail = generate_thumbnail(slide_pptx_path, thumbnails_dir / thumbnail_filename, content_bbox_pt)
            thumbnail_rel = f"{category}/{thumbnail_filename}" if has_thumbnail else None

            items.append(
                {
                    "title": title,
                    "insertMode": "file",
                    "sourceFile": f"{category}/{slide_filename}",
                    "thumbnail": thumbnail_rel,
                    "sortOrder": item_number,
                }
            )

    seed_path = Path("db/seed") / f"catalog-{category}.json"
    seed_path.parent.mkdir(parents=True, exist_ok=True)
    seed_path.write_text(json.dumps({"category": category, "items": items}, indent=2) + "\n")

    print(f"Wrote {len(items)} item(s) to {seed_path}", file=sys.stderr)
    print(f"Sliced {len(items)} 'file'-mode .pptx file(s) into: {category_dir}", file=sys.stderr)
    print(f"Generated thumbnails into: {thumbnails_dir}", file=sys.stderr)
    print(
        f"Ready to seed as-is: node scripts/seed-catalog.js {seed_path}\n"
        f"Titles are the country names as found - correct via /admin after seeding, not by re-running this script.",
        file=sys.stderr,
    )


if __name__ == "__main__":
    main()
