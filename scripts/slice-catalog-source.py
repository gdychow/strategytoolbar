#!/usr/bin/env python3
"""
Content-prep step for Tier 3's shared catalog. For a given boilerplate
library .pptx (e.g. Text.pptx), classifies each slide's real content (the
inert think-cell placeholder frame is always stripped first, regardless of
mode) as either:

  reconstruct - every real shape is a plain text box, or preset geometry
    (<a:prstGeom>) whose prst has a confirmed, hand-checked entry in
    PRST_TO_GEOMETRIC_SHAPE_TYPE below - so it can be rebuilt with full
    fidelity via addGeometricShape/addTextBox/addGroup at insert time (see
    src/features/libraryInsert.ts). No .pptx file is kept for these items
    long-term - one is sliced to a scratch directory purely to render its
    thumbnail, then discarded.

  file - anything else (custom <a:custGeom> geometry, embedded pictures,
    tables/other graphicFrames, or a prstGeom whose preset isn't in
    PRST_TO_GEOMETRIC_SHAPE_TYPE) - no PowerPoint JS API can reconstruct
    the first group, and the last group is deliberately not guessed at:
    addGeometricShape would silently accept a wrong-but-validly-spelled
    enum value and render the wrong shape, with no error. Either way the
    slide is sliced out into its own minimal single-slide .pptx under
    <output-dir>/<category>/, to be inserted via insertSlidesFromBase64 at
    runtime.

Every item (both modes) also gets a thumbnail, rendered via macOS
QuickLook (`qlmanage -t`, a stock command-line tool - no LibreOffice/
poppler install needed, confirmed against a real sliced file) into
<output-dir>/thumbnails/<category>/, and a rough title auto-filled from
the slide's own text (or "<Category> #N" for text-less shapes/graphics).
Titles and thumbnails generated this way are meant to be rough starting
points, not final - db/seed/catalog-<category>.json is written directly
in the exact shape scripts/seed-catalog.js expects, ready to seed
immediately; use /admin afterward to correct any title, thumbnail, or
category assignment that needs a human's judgment, without re-running
this script or redeploying.

Thumbnails are cropped to the real content shapes' combined bounding box
(plus padding), not the whole 16x9 slide - most catalog items are small
relative to a full slide canvas, and a full-slide thumbnail is mostly
white space. Rendered at a higher resolution than the final crop needs
(see RENDER_SIZE) so a small, tightly-cropped shape still has enough
source pixels to look sharp rather than blurry.

Usage:
    python3 scripts/slice-catalog-source.py <source.pptx> <category-slug> [output-dir]

Example:
    python3 scripts/slice-catalog-source.py \
        "../Package Files/StrategyToolbar/BoilerPlates/16x9/Objects.pptx" \
        objects data/catalog
"""
import json
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.util import Emu

THINK_CELL_MARKER = "think-cell"
EMU_PER_POINT = 12700
# Every category built so far targets this width. Some older boilerplate
# decks (Maps, ClipArt, Frameworks, Flags) are 720x540pt (4:3) instead -
# confirmed their height already matches exactly (540pt), so centering
# only ever needs a horizontal shift, never rescaling. See
# TARGET_SLIDE_WIDTH_PT's use in main()/slice_single_slide/
# extract_reconstruct_spec.
TARGET_SLIDE_WIDTH_PT = 960.0
# Whole-slide render resolution fed into the content-bbox crop below - much
# higher than the ~320px final thumbnail size, since a small shape crops
# down to a small fraction of this; too low a render size and small items
# end up visibly blurry once cropped.
RENDER_SIZE = 1600
CONTENT_PADDING_FRAC = 0.15  # of the content bbox's larger dimension
CONTENT_PADDING_MIN_PT = 8


def is_think_cell_placeholder(shape) -> bool:
    return THINK_CELL_MARKER in shape._element.xml


def emu_to_pt(value) -> float:
    return round(Emu(value).pt, 2) if value is not None else None


def has_cust_geom(element) -> bool:
    return element.find(".//" + qn("a:custGeom")) is not None


def has_picture(element) -> bool:
    # element.tag check needed alongside .//: a top-level (non-grouped)
    # picture shape's own _element IS the <p:pic> node - "shape._element.find('.//p:pic')"
    # only searches descendants, so it misses a standalone picture entirely
    # (found via Symbols.pptx slides 6-7, each a single top-level <p:pic>
    # with a plain prstGeom="rect" bounding box - the missed check let
    # classify_shape_tree fall through to the prstGeom mapping and
    # misclassify these as a reconstructable empty rectangle, silently
    # dropping the actual image). Nested pictures (inside a group) are
    # still correctly caught by the descendant search alone.
    return element.tag == qn("p:pic") or element.find(".//" + qn("p:pic")) is not None


def has_other_graphic_frame(element) -> bool:
    # Any graphicFrame other than the (already-excluded) think-cell one -
    # tables, charts, SmartArt. Not handled by the reconstruct path.
    # Same self-vs-descendant gap as has_picture() above: a top-level
    # graphicFrame's own _element IS the <p:graphicFrame> node, which a
    # pure ".//" descendant search misses. Not observed to matter in
    # practice yet (a graphicFrame has no prstGeom, so classify_shape_tree
    # already falls through to 'file' mode via the "no mapped preset"
    # branch either way) - fixed anyway for correctness, not defensively.
    return element.tag == qn("p:graphicFrame") or element.find(".//" + qn("p:graphicFrame")) is not None


def has_custom_bullet_char(element) -> bool:
    # PowerPoint.BulletFormat (confirmed in src/features/otherTweaks.ts,
    # while porting L_Other_Tweaks.bas) only exposes type
    # (None/Numbered/Unnumbered), a fixed numbering-style enum, and
    # visible - no way to set a specific bullet character, font, or
    # colour. A <a:buChar> (an explicit fixed bullet glyph, as opposed to
    # <a:buAutoNum>'s auto-incrementing numbering) can't be reproduced, so
    # route it to the file-based path instead of guessing at a lossy
    # approximation.
    return element.find(".//" + qn("a:buChar")) is not None


# Maps an OOXML <a:prstGeom prst="..."> value to the exact string
# PowerPoint.GeometricShapeType's TypeScript enum uses (per
# node_modules/@types/office-js), which src/features/libraryInsert.ts
# passes straight to addGeometricShape with no translation of its own.
# These are NOT simply the prst value capitalized - some are different
# words entirely ("rect" -> "Rectangle", not "Rect") - so this is a
# deliberately conservative, hand-checked subset against the real enum
# list, not a guessed/mechanical transform. A prst NOT in this table
# routes to 'file' mode instead (see classify_shape_tree) rather than
# risk emitting a wrong-but-validly-spelled enum value, which
# addGeometricShape would silently accept and render as the WRONG shape
# with no error - a much worse failure mode than falling back to the
# already-proven file-mode insert path.
PRST_TO_GEOMETRIC_SHAPE_TYPE = {
    "rect": "Rectangle",
    "ellipse": "Ellipse",
    "roundRect": "RoundRectangle",
    "triangle": "Triangle",
    "rtTriangle": "RightTriangle",
    "diamond": "Diamond",
    "parallelogram": "Parallelogram",
    "trapezoid": "Trapezoid",
    "pentagon": "Pentagon",
    "hexagon": "Hexagon",
    "heptagon": "Heptagon",
    "octagon": "Octagon",
    "decagon": "Decagon",
    "dodecagon": "Dodecagon",
    "star4": "Star4",
    "star5": "Star5",
    "star6": "Star6",
    "star7": "Star7",
    "star8": "Star8",
    "star10": "Star10",
    "star12": "Star12",
    "star16": "Star16",
    "star24": "Star24",
    "star32": "Star32",
    "chevron": "Chevron",
    "donut": "Donut",
    "heart": "Heart",
    "sun": "Sun",
    "moon": "Moon",
    "cube": "Cube",
    "can": "Can",
    "cloud": "Cloud",
    "wave": "Wave",
    "plus": "Plus",
    "arc": "Arc",
    "chord": "Chord",
    "bevel": "Bevel",
    "frame": "Frame",
    "corner": "Corner",
    "pie": "Pie",
    "blockArc": "BlockArc",
    "smileyFace": "SmileyFace",
    "lightningBolt": "LightningBolt",
    "plaque": "Plaque",
    "teardrop": "Teardrop",
    "homePlate": "HomePlate",
    "rightArrow": "RightArrow",
    "leftArrow": "LeftArrow",
    "upArrow": "UpArrow",
    "downArrow": "DownArrow",
    "leftRightArrow": "LeftRightArrow",
    "upDownArrow": "UpDownArrow",
    "leftBracket": "LeftBracket",
    "rightBracket": "RightBracket",
    "leftBrace": "LeftBrace",
    "rightBrace": "RightBrace",
    "wedgeRectCallout": "WedgeRectCallout",
    # Deliberately NOT mapped: prst="line" has no GeometricShapeType
    # equivalent at all (confirmed against the real addGeometricShape
    # method signature's full literal-union parameter type, not just the
    # enum declaration) - a straight line needs PowerPoint.Slide.addLine,
    # a different API this project doesn't call from the reconstruct path
    # yet. Left unmapped so classify_shape_tree correctly routes it to
    # file mode instead.
}


def classify_shape_tree(shape) -> str:
    """Returns 'file' if anything under this shape can't be reconstructed, else 'reconstruct'."""
    el = shape._element
    if has_cust_geom(el) or has_picture(el) or has_other_graphic_frame(el) or has_custom_bullet_char(el):
        return "file"
    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        # extract_reconstruct_spec never reads presetGeometry for a text
        # box (it builds it via addTextBox, not addGeometricShape), so its
        # prstGeom - if any - doesn't matter here.
        return "reconstruct"
    # Any other shape needs a prstGeom with a confirmed, mapped preset -
    # covers both "no prstGeom at all" (e.g. a connector shape) and "a
    # prst we don't have a verified enum value for" (see
    # PRST_TO_GEOMETRIC_SHAPE_TYPE's comment for why an unmapped preset
    # falls back to file mode rather than guessing).
    prst_el = el.find(".//" + qn("a:prstGeom"))
    prst = prst_el.get("prst") if prst_el is not None else None
    if prst not in PRST_TO_GEOMETRIC_SHAPE_TYPE:
        return "file"
    return "reconstruct"


def extract_fill(shape):
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        if str(fill.type) == "MSO_FILL_TYPE.SOLID (1)" or fill.type == 1:
            return {"type": "solid", "color": f"#{fill.fore_color.rgb}"}
        return {"type": str(fill.type), "note": "non-solid fill - verify manually against ShapeFill API"}
    except Exception:
        return None


def extract_line(shape):
    try:
        line = shape.line
        if line.fill.type is None:
            return None
        width_pt = emu_to_pt(line.width) if line.width else None
        try:
            color = f"#{line.color.rgb}"
        except Exception:
            color = None
        return {"color": color, "widthPt": width_pt}
    except Exception:
        return None


def extract_bullet(paragraph_xml):
    bu_char = paragraph_xml.find(".//" + qn("a:buChar"))
    if bu_char is not None:
        return bu_char.get("char")
    if paragraph_xml.find(".//" + qn("a:buNone")) is not None:
        return None
    return "(inherited - verify against shape's lstStyle or slide layout)"


def extract_text_spec(shape):
    if not shape.has_text_frame:
        return None
    paragraphs = []
    for p in shape.text_frame.paragraphs:
        runs = []
        for r in p.runs:
            runs.append(
                {
                    "text": r.text,
                    "bold": r.font.bold,
                    "italic": r.font.italic,
                    "size": r.font.size.pt if r.font.size else None,
                    "fontName": r.font.name,
                    "color": (f"#{r.font.color.rgb}" if r.font.color and r.font.color.type is not None else None),
                }
            )
        paragraphs.append({"level": p.level, "bullet": extract_bullet(p._p), "runs": runs})
    return paragraphs


def extract_reconstruct_spec(shape, x_offset_pt: float = 0.0):
    is_text_box = shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    spec = {
        "kind": "textBox" if is_text_box else "geometricShape",
        "left": round(emu_to_pt(shape.left) + x_offset_pt, 2),
        "top": emu_to_pt(shape.top),
        "width": emu_to_pt(shape.width),
        "height": emu_to_pt(shape.height),
        "rotation": shape.rotation,
        "fill": extract_fill(shape),
        "line": extract_line(shape),
        # Must match ShapeSpec.paragraphs in src/features/libraryInsert.ts,
        # which is what insertReconstructedItem's applyParagraphs() actually
        # reads - this key was previously named "text", a silent mismatch
        # that made every script-seeded reconstruct-mode item's real text
        # content get dropped on insert with no error (found while curating
        # the Symbols category; confirmed against catalog-objects/stamps/
        # diagrams.json, which shipped with the wrong key).
        "paragraphs": extract_text_spec(shape),
    }
    if not is_text_box:
        prst_el = shape._element.find(".//" + qn("a:prstGeom"))
        prst = prst_el.get("prst") if prst_el is not None else None
        # classify_shape_tree() only classifies a shape 'reconstruct' when
        # its prst is a confirmed key in PRST_TO_GEOMETRIC_SHAPE_TYPE, so
        # this is a direct (not .get-with-fallback) lookup on purpose - a
        # KeyError here means that invariant broke, which should fail
        # loudly rather than silently emit an unmapped/wrong value.
        spec["presetGeometry"] = PRST_TO_GEOMETRIC_SHAPE_TYPE[prst]
    return spec


def extract_title_from_placeholder(shapes):
    """
    Prefers a real Title/Center-Title placeholder's own text over the
    generic multi-shape text-hint heuristic below. Text/Objects/Shapes/etc.
    have no such placeholder (hence the hint heuristic existing at all),
    but Maps/ClipArt/Frameworks/Flags do - confirmed by inspecting real
    slides directly - so this gives a real name ("Abu Dhabi Island",
    "Africa A-G") instead of a rough concatenation of whatever text happens
    to be on the slide. Returns None (falls back to the hint) when no
    title placeholder is present.
    """
    for shape in shapes:
        if not shape.is_placeholder:
            continue
        ph_type = shape.placeholder_format.type
        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE) and shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text
    return None


def extract_text_hint(shape) -> str:
    if not shape.has_text_frame:
        return ""
    # p.text preserves "\x0b" for a <a:br/> soft line break within a
    # paragraph (e.g. "Total Savings\x0b= $23 Billion") - a literal
    # control character, not just cosmetically odd, so it's replaced
    # rather than left for a human to notice later in /admin.
    texts = (p.text.replace("\x0b", " ").strip() for p in shape.text_frame.paragraphs)
    return " / ".join(t for t in texts if t)[:80]


def slice_single_slide(
    source_path: Path, slide_index: int, dest_path: Path, x_offset_pt: float = 0.0, target_width_pt: float = None
) -> None:
    """
    Duplicates the source presentation, keeps only slide_index, strips the
    think-cell placeholder, saves to dest_path. x_offset_pt/target_width_pt
    are only passed for 'file'-mode items from a non-16x9 source (see
    TARGET_SLIDE_WIDTH_PT) - shifts every remaining shape right by
    x_offset_pt and resizes the presentation itself to target_width_pt, so
    the packaged single-slide file is already correctly centered for a
    16x9 destination rather than relying on insertSlidesFromBase64's own
    (unverified) handling of a slide-size mismatch. Left at their defaults
    (no-op) for the scratch thumbnail-only slice of a 'reconstruct'-mode
    item, since that render only needs to match the *original* slide's own
    coordinate space that compute_content_bbox_pt already uses.
    """
    prs = Presentation(str(source_path))
    slide_id_list = prs.slides._sldIdLst
    slide_id_elements = list(slide_id_list)
    for i, slide_id_elm in enumerate(slide_id_elements):
        if i == slide_index:
            continue
        rId = slide_id_elm.get(qn("r:id"))
        prs.part.drop_rel(rId)
        slide_id_list.remove(slide_id_elm)

    remaining_slide = prs.slides[0]
    for shape in list(remaining_slide.shapes):
        if is_think_cell_placeholder(shape):
            shape._element.getparent().remove(shape._element)

    if x_offset_pt:
        offset_emu = round(x_offset_pt * EMU_PER_POINT)
        for shape in remaining_slide.shapes:
            shape.left = shape.left + offset_emu
    if target_width_pt is not None:
        prs.slide_width = round(target_width_pt * EMU_PER_POINT)

    dest_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest_path))


def compute_content_bbox_pt(shapes, slide_width_pt: float, slide_height_pt: float):
    """
    Combined bounding box (left, top, right, bottom, in points) of every
    given shape, padded by CONTENT_PADDING_FRAC of its own larger
    dimension (at least CONTENT_PADDING_MIN_PT) and clamped to the slide's
    canvas - used to crop a full-slide thumbnail render down to the actual
    content instead of a mostly-blank 16x9 slide.
    """
    lefts, tops, rights, bottoms = [], [], [], []
    for shape in shapes:
        left, top = emu_to_pt(shape.left), emu_to_pt(shape.top)
        width, height = emu_to_pt(shape.width), emu_to_pt(shape.height)
        lefts.append(left)
        tops.append(top)
        rights.append(left + width)
        bottoms.append(top + height)

    left, top, right, bottom = min(lefts), min(tops), max(rights), max(bottoms)
    pad = max(CONTENT_PADDING_FRAC * max(right - left, bottom - top), CONTENT_PADDING_MIN_PT)
    left, top = max(0.0, left - pad), max(0.0, top - pad)
    right, bottom = min(slide_width_pt, right + pad), min(slide_height_pt, bottom + pad)
    return left, top, right, bottom


def generate_thumbnail(slide_pptx_path: Path, dest_path: Path, content_bbox_pt, slide_width_pt: float, slide_height_pt: float) -> bool:
    """
    Renders slide_pptx_path (a single-slide .pptx) to a PNG via macOS
    QuickLook - no LibreOffice/poppler install needed - then crops to
    content_bbox_pt (see compute_content_bbox_pt). Returns False (and logs
    a warning) if QuickLook doesn't produce output for this slide -
    observed occasionally for unusual content - so the caller can skip the
    thumbnail rather than aborting the whole category run.
    """
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            ["qlmanage", "-t", "-s", str(RENDER_SIZE), "-o", tmp, str(slide_pptx_path)],
            capture_output=True,
            check=False,
        )
        produced = Path(tmp) / f"{slide_pptx_path.name}.png"
        if not produced.exists():
            print(f"  warning: qlmanage produced no thumbnail for {slide_pptx_path.name}", file=sys.stderr)
            return False

        dest_path.parent.mkdir(parents=True, exist_ok=True)
        img = Image.open(produced)
        scale_x = img.width / slide_width_pt
        scale_y = img.height / slide_height_pt
        left, top, right, bottom = content_bbox_pt
        box_px = (round(left * scale_x), round(top * scale_y), round(right * scale_x), round(bottom * scale_y))
        img.crop(box_px).save(dest_path)
        return True


def main():
    args = sys.argv[1:]
    # --group is only used for Maps (four separate source files - MapsI-IV -
    # feeding into one "maps" category): tags every item this invocation
    # produces with that group name, and - since the seed file will already
    # exist from an earlier MapsX run - appends to it (continuing sortOrder/
    # filename numbering from its current max) instead of overwriting.
    # Every other, single-invocation category never passes this, so it
    # keeps today's exact behavior (fresh overwrite, no group).
    group_name = None
    if "--group" in args:
        idx = args.index("--group")
        group_name = args[idx + 1]
        del args[idx : idx + 2]

    if len(args) < 2:
        print(__doc__)
        sys.exit(1)

    source_path = Path(args[0])
    category = args[1]
    output_root = Path(args[2]) if len(args) > 2 else Path("data/catalog")
    category_dir = output_root / category
    thumbnails_dir = output_root / "thumbnails" / category

    prs = Presentation(str(source_path))
    slide_width_pt, slide_height_pt = emu_to_pt(prs.slide_width), emu_to_pt(prs.slide_height)
    # 0.0 for already-16x9 sources (every category before this one) - only
    # the older 4:3 decks (Maps/ClipArt/Frameworks/Flags) get a real shift.
    x_offset_pt = round((TARGET_SLIDE_WIDTH_PT - slide_width_pt) / 2, 2)

    seed_path = Path("db/seed") / f"catalog-{category}.json"
    existing_items = []
    starting_number = 0
    if group_name is not None and seed_path.exists():
        existing = json.loads(seed_path.read_text())
        existing_items = existing.get("items", [])
        starting_number = max((item.get("sortOrder", 0) for item in existing_items), default=0)

    items = []
    running_item_count = 0

    with tempfile.TemporaryDirectory() as scratch:
        scratch_dir = Path(scratch)

        for index, slide in enumerate(prs.slides):
            real_shapes = [s for s in slide.shapes if not is_think_cell_placeholder(s)]
            if not real_shapes:
                print(f"slide {index + 1}: no real content shape found, skipping", file=sys.stderr)
                continue

            # Append-mode numbering must continue across separate
            # invocations (MapsI/II/III/IV each restart their own slide
            # index at 1, which would otherwise collide); every other
            # category keeps today's exact per-slide-index numbering.
            if group_name is not None:
                running_item_count += 1
                item_number = starting_number + running_item_count
            else:
                item_number = index + 1

            mode = "reconstruct" if all(classify_shape_tree(s) == "reconstruct" for s in real_shapes) else "file"
            title = extract_title_from_placeholder(real_shapes)
            if not title:
                text_hint = " | ".join(filter(None, (extract_text_hint(s) for s in real_shapes)))
                title = text_hint if text_hint else f"{category.capitalize()} #{item_number}"
            slide_filename = f"{category}-{item_number:03d}.pptx"

            # Every item gets sliced to a single-slide .pptx, regardless of
            # mode - for 'file' items this is the real, permanent catalog
            # content; for 'reconstruct' items it's a scratch file that
            # exists only long enough to render a faithful thumbnail from
            # (the real content is the reconstructSpec JSON below). Only
            # 'file' items get the 4:3->16x9 offset/resize - the scratch
            # slice's thumbnail render doesn't need it (see
            # slice_single_slide's docstring).
            if mode == "file":
                slide_pptx_path = category_dir / slide_filename
                slice_single_slide(source_path, index, slide_pptx_path, x_offset_pt, TARGET_SLIDE_WIDTH_PT)
            else:
                slide_pptx_path = scratch_dir / slide_filename
                slice_single_slide(source_path, index, slide_pptx_path)

            content_bbox_pt = compute_content_bbox_pt(real_shapes, slide_width_pt, slide_height_pt)
            thumbnail_filename = f"{category}-{item_number:03d}.png"
            has_thumbnail = generate_thumbnail(
                slide_pptx_path, thumbnails_dir / thumbnail_filename, content_bbox_pt, slide_width_pt, slide_height_pt
            )
            thumbnail_rel = f"{category}/{thumbnail_filename}" if has_thumbnail else None

            if mode == "file":
                item = {
                    "title": title,
                    "insertMode": "file",
                    "sourceFile": f"{category}/{slide_filename}",
                    "thumbnail": thumbnail_rel,
                    "sortOrder": item_number,
                }
            else:
                if len(real_shapes) == 1:
                    spec = extract_reconstruct_spec(real_shapes[0], x_offset_pt)
                else:
                    spec = {"kind": "group", "shapes": [extract_reconstruct_spec(s, x_offset_pt) for s in real_shapes]}
                item = {
                    "title": title,
                    "insertMode": "reconstruct",
                    "reconstructSpec": spec,
                    "thumbnail": thumbnail_rel,
                    "sortOrder": item_number,
                }
            if group_name is not None:
                item["groupName"] = group_name
            items.append(item)

    seed_path.parent.mkdir(parents=True, exist_ok=True)
    all_items = existing_items + items
    seed_path.write_text(json.dumps({"category": category, "items": all_items}, indent=2) + "\n")

    file_mode_count = sum(1 for i in items if i["insertMode"] == "file")
    print(f"Wrote {len(items)} item(s) ({len(all_items)} total in file) to {seed_path}", file=sys.stderr)
    if file_mode_count:
        print(f"Sliced {file_mode_count} 'file'-mode .pptx file(s) into: {category_dir}", file=sys.stderr)
    print(f"Generated thumbnails into: {thumbnails_dir}", file=sys.stderr)
    print(
        f'Ready to seed as-is: node scripts/seed-catalog.js {seed_path}\n'
        f"Titles are rough (auto-filled from slide text) - correct via /admin after seeding, not by re-running this script.",
        file=sys.stderr,
    )


if __name__ == "__main__":
    main()
