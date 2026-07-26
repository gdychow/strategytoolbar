#!/usr/bin/env python3
"""
One-off repair script for a real bug found in real PowerPoint: slice-
catalog-source.py's 4:3->16x9 offset fix shifted every shape's .left on a
sliced 'file'-mode slide, including placeholder shapes (e.g. a title) that
had no explicit <a:xfrm> of their own - they were relying on inheriting
position/size from their slide layout. Setting .left alone on such a shape
makes python-pptx write a *new* xfrm containing only <a:off> (position),
never a <a:ext> (size) - confirmed directly. PowerPoint then renders that
shape at an effectively zero width instead of falling back to the
inherited size (observed: a title's text wrapping to one character per
line). Only affects Maps/ClipArt/Frameworks (the 4:3-sourced categories -
Flags builds brand-new slides with no inherited placeholders at all, so
it was never affected).

Fix: re-set every placeholder shape's width/height explicitly (reading
its own current, already-correctly-inherited value via python-pptx's
getter, then writing it straight back) - this materializes a complete
xfrm with both <a:off> and <a:ext>, regardless of whether one existed
before. A no-op for any shape that already had a complete xfrm.

Repairs already-sliced .pptx files in place - no need to re-run the
original slicing (source shapes/positions are already correct; only the
missing <a:ext> needs fixing), so no thumbnail regeneration, no seed
JSON change, and no reseed needed - just overwrite the existing files at
their existing paths.

Usage:
    python3 scripts/fix-placeholder-xfrm.py data/catalog/maps data/catalog/clipart data/catalog/frameworks
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn


def has_complete_xfrm(shape) -> bool:
    spPr = shape._element.spPr
    if spPr is None:
        return False
    xfrm = spPr.find(qn("a:xfrm"))
    return xfrm is not None and xfrm.find(qn("a:ext")) is not None


def repair_file(path: Path) -> bool:
    """Returns True if the file needed (and got) a repair."""
    prs = Presentation(str(path))
    changed = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.is_placeholder or has_complete_xfrm(shape):
                continue
            width, height = shape.width, shape.height
            shape.width = width
            shape.height = height
            changed = True
    if changed:
        prs.save(str(path))
    return changed


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    total_files = 0
    total_repaired = 0
    for dir_arg in sys.argv[1:]:
        directory = Path(dir_arg)
        for pptx_path in sorted(directory.glob("*.pptx")):
            total_files += 1
            if repair_file(pptx_path):
                total_repaired += 1
                print(f"repaired: {pptx_path}", file=sys.stderr)

    print(f"Checked {total_files} file(s), repaired {total_repaired}.", file=sys.stderr)


if __name__ == "__main__":
    main()
